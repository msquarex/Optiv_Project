"""
File Processing Module for VIT Campus Connect
Handles multiple file formats and extracts text content
"""

import os
import pandas as pd
import PyPDF2
import pytesseract
from PIL import Image
import cv2
import numpy as np
from docx import Document
from pptx import Presentation
import logging
from typing import Dict, List, Tuple, Any
import io
import fitz  # PyMuPDF
import tempfile
import re
from image_descriptor import ImageDescriptor

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class FileProcessor:
    """Main file processing class for handling multiple file formats"""
    
    def __init__(self, gemma_url: str = "http://127.0.0.1:11500"):
        self.supported_formats = {
            '.xlsx': self._process_excel,
            '.xls': self._process_excel,
            '.csv': self._process_csv,
            '.txt': self._process_text,
            '.pdf': self._process_pdf,
            '.jpeg': self._process_image,
            '.jpg': self._process_image,
            '.png': self._process_image,
            '.pptx': self._process_powerpoint,
            '.docx': self._process_word
        }
        # Initialize image descriptor for Gemma3 integration
        self.image_descriptor = ImageDescriptor(gemma_url)
    
    def process_file(self, file_path: str, enable_image_descriptions: bool = True) -> Dict[str, Any]:
        """
        Main method to process any supported file format
        
        Args:
            file_path (str): Path to the file to process
            enable_image_descriptions (bool): Whether to enable AI image descriptions for image files
            
        Returns:
            Dict containing extracted text and metadata
        """
        try:
            file_ext = os.path.splitext(file_path)[1].lower()
            
            if file_ext not in self.supported_formats:
                raise ValueError(f"Unsupported file format: {file_ext}")
            
            logger.info(f"Processing file: {file_path}")
            if file_ext in ['.png', '.jpg', '.jpeg']:
                result = self.supported_formats[file_ext](file_path, enable_image_descriptions=enable_image_descriptions)
            else:
                result = self.supported_formats[file_ext](file_path)
            
            # Handle different result formats
            if isinstance(result, dict):
                # Image processing returns a dict with multiple fields
                if 'combined_content' in result:
                    # This is an image with both OCR and description
                    return {
                        'file_path': file_path,
                        'file_type': file_ext,
                        'content': result['combined_content'],
                        'ocr_text': result.get('ocr_text', ''),
                        'image_description': result.get('image_description', ''),
                        'description_status': result.get('description_status', ''),
                        'description_error': result.get('description_error', ''),
                        'status': result.get('status', 'success')
                    }
                else:
                    # Other dict results
                    return {
                        'file_path': file_path,
                        'file_type': file_ext,
                        'content': result.get('content', str(result)),
                        'status': result.get('status', 'success')
                    }
            else:
                # String results (most other file types)
                return {
                    'file_path': file_path,
                    'file_type': file_ext,
                    'content': result,
                    'status': 'success'
                }
            
        except Exception as e:
            logger.error(f"Error processing file {file_path}: {str(e)}")
            return {
                'file_path': file_path,
                'file_type': file_ext if 'file_ext' in locals() else 'unknown',
                'content': '',
                'status': 'error',
                'error': str(e)
            }
    
    def _process_excel(self, file_path: str) -> str:
        """Process Excel files and extract text content"""
        try:
            # Read all sheets
            excel_file = pd.ExcelFile(file_path)
            content_parts = []
            
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                content_parts.append(f"Sheet: {sheet_name}")
                content_parts.append(df.to_string())
                content_parts.append("\n" + "="*50 + "\n")
            
            return "\n".join(content_parts)
            
        except Exception as e:
            logger.error(f"Error processing Excel file: {str(e)}")
            return f"Error reading Excel file: {str(e)}"
    
    def _process_csv(self, file_path: str) -> str:
        """Process CSV files and extract text content"""
        try:
            df = pd.read_csv(file_path)
            return df.to_string()
        except Exception as e:
            logger.error(f"Error processing CSV file: {str(e)}")
            return f"Error reading CSV file: {str(e)}"
    
    def _process_text(self, file_path: str) -> str:
        """Process plain text files"""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
                return file.read()
        except Exception as e:
            logger.error(f"Error processing text file: {str(e)}")
            return f"Error reading text file: {str(e)}"
    
    def _process_pdf(self, file_path: str) -> str:
        """Process PDF files using PyPDF2 and OCR for scanned PDFs"""
        try:
            content_parts = []
            pages_needing_ocr = []
            
            # First pass: Try to extract text directly
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                total_pages = len(pdf_reader.pages)
                
                for page_num, page in enumerate(pdf_reader.pages):
                    try:
                        text = page.extract_text()
                        if text.strip():
                            content_parts.append(f"Page {page_num + 1}:")
                            content_parts.append(text)
                        else:
                            # Mark page for OCR processing
                            pages_needing_ocr.append(page_num)
                            content_parts.append(f"Page {page_num + 1}: [Scanned content - OCR processing needed]")
                    except Exception as e:
                        logger.warning(f"Error extracting text from page {page_num + 1}: {str(e)}")
                        pages_needing_ocr.append(page_num)
                        content_parts.append(f"Page {page_num + 1}: [Error extracting content]")
            
            # Second pass: Process scanned pages with OCR
            if pages_needing_ocr:
                logger.info(f"Processing {len(pages_needing_ocr)} pages with OCR")
                ocr_results = self._process_pdf_ocr(file_path, pages_needing_ocr)
                
                # Replace placeholder text with OCR results
                for i, page_num in enumerate(pages_needing_ocr):
                    placeholder = f"Page {page_num + 1}: [Scanned content - OCR processing needed]"
                    error_placeholder = f"Page {page_num + 1}: [Error extracting content]"
                    
                    # Find and replace the placeholder
                    for j, content in enumerate(content_parts):
                        if placeholder in content or error_placeholder in content:
                            content_parts[j] = f"Page {page_num + 1}:"
                            # Insert OCR result after the page header
                            content_parts.insert(j + 1, ocr_results[i])
                            break
            
            return "\n".join(content_parts)
            
        except Exception as e:
            logger.error(f"Error processing PDF file: {str(e)}")
            return f"Error reading PDF file: {str(e)}"
    
    def _process_pdf_ocr(self, file_path: str, page_numbers: List[int]) -> List[str]:
        """Process specific PDF pages with OCR using PyMuPDF"""
        try:
            # Open PDF with PyMuPDF
            pdf_document = fitz.open(file_path)
            ocr_results = []
            
            for page_num in page_numbers:
                try:
                    # Get the page (PyMuPDF uses 0-based indexing)
                    page = pdf_document[page_num]
                    
                    # Convert page to image with high resolution
                    mat = fitz.Matrix(2.0, 2.0)  # 2x zoom for better OCR quality
                    pix = page.get_pixmap(matrix=mat)
                    
                    # Convert to PIL Image
                    img_data = pix.tobytes("png")
                    image = Image.open(io.BytesIO(img_data))
                    
                    # Convert PIL image to OpenCV format
                    opencv_image = cv2.cvtColor(np.array(image), cv2.COLOR_RGB2BGR)
                    
                    # Preprocess image for better OCR
                    gray = cv2.cvtColor(opencv_image, cv2.COLOR_BGR2GRAY)
                    
                    # Apply denoising
                    denoised = cv2.fastNlMeansDenoising(gray)
                    
                    # Apply thresholding for better text recognition
                    _, thresh = cv2.threshold(denoised, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)
                    
                    # Use pytesseract for OCR with optimized settings
                    custom_config = r'--oem 3 --psm 6'
                    text = pytesseract.image_to_string(thresh, config=custom_config)
                    
                    # Clean up the text
                    cleaned_text = self._clean_ocr_text(text)
                    ocr_results.append(cleaned_text if cleaned_text.strip() else "[No text detected via OCR]")
                    
                    logger.info(f"OCR completed for page {page_num + 1}")
                    
                except Exception as e:
                    logger.error(f"OCR error for page {page_num + 1}: {str(e)}")
                    ocr_results.append(f"[OCR processing failed: {str(e)}]")
            
            # Close the PDF document
            pdf_document.close()
            return ocr_results
            
        except Exception as e:
            logger.error(f"Error in PDF OCR processing: {str(e)}")
            return [f"[OCR processing failed: {str(e)}]" for _ in page_numbers]
    
    def _clean_ocr_text(self, text: str) -> str:
        """Clean and format OCR text output"""
        if not text:
            return ""
        
        # Remove excessive whitespace and normalize line breaks
        lines = text.split('\n')
        cleaned_lines = []
        
        for line in lines:
            # Remove leading/trailing whitespace
            cleaned_line = line.strip()
            
            # Skip empty lines
            if cleaned_line:
                cleaned_lines.append(cleaned_line)
        
        # Join lines with proper spacing
        return '\n'.join(cleaned_lines)
    
    def _ocr_simple_methods(self, gray_image: np.ndarray) -> str:
        """Simple OCR methods that work best for clean, simple text"""
        try:
            results = []
            
            # Try PSM 3 (fully automatic page segmentation, but no OSD)
            try:
                text = pytesseract.image_to_string(gray_image, config='--oem 3 --psm 3')
                cleaned = self._clean_ocr_text(text)
                if cleaned.strip():
                    results.append(cleaned)
            except Exception as e:
                logger.warning(f"PSM 3 failed: {str(e)}")
            
            # Try PSM 4 (assume a single column of text of variable sizes)
            try:
                text = pytesseract.image_to_string(gray_image, config='--oem 3 --psm 4')
                cleaned = self._clean_ocr_text(text)
                if cleaned.strip():
                    results.append(cleaned)
            except Exception as e:
                logger.warning(f"PSM 4 failed: {str(e)}")
            
            # Return the cleanest result (prefer shorter, cleaner results)
            if results:
                # If we have multiple results, prefer the one that looks most like clean text
                # (shorter length often indicates less noise)
                return min(results, key=lambda x: len(x.strip()))
            
            return ""
            
        except Exception as e:
            logger.error(f"Error in simple OCR methods: {str(e)}")
            return ""
    
    def _process_image(self, file_path: str, enable_image_descriptions: bool = True) -> Dict[str, Any]:
        """Process image files using OCR and Gemma3 image description"""
        try:
            # Load image
            image = cv2.imread(file_path)
            if image is None:
                return {
                    'ocr_text': f"Error: Could not load image from {file_path}",
                    'image_description': '',
                    'status': 'error',
                    'error': f"Could not load image from {file_path}"
                }
            
            # Convert to grayscale for better OCR performance
            gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
            
            # Try multiple OCR approaches and combine results
            ocr_results = []
            
            # Approach 1: Simple grayscale with PSM 3 and 4 (most accurate for clean text)
            ocr_results.append(self._ocr_simple_methods(gray))
            
            # Approach 2: Enhanced preprocessing only if simple methods fail
            simple_result = self._ocr_simple_methods(gray)
            if len(simple_result.strip()) < 5:  # If simple method gives very little text
                ocr_results.append(self._ocr_with_preprocessing(image, 'standard'))
                ocr_results.append(self._ocr_with_preprocessing(image, 'handwritten'))
                ocr_results.append(self._ocr_with_preprocessing(image, 'table'))
            
            # Approach 3: Raw image with different PSM modes as fallback
            if not any(result.strip() for result in ocr_results):
                ocr_results.append(self._ocr_raw_image(image))
            
            # Combine and clean OCR results
            combined_ocr_text = self._combine_ocr_results(ocr_results)
            final_ocr_text = combined_ocr_text if combined_ocr_text.strip() else "No text detected in image"
            
            # Get image description from Gemma3 (if enabled)
            if enable_image_descriptions:
                logger.info(f"Getting image description from Gemma3 for: {file_path}")
                description_result = self.image_descriptor.describe_image(file_path)
            else:
                logger.info(f"Skipping image description for: {file_path} (disabled by user)")
                description_result = {
                    'status': 'skipped',
                    'description': '',
                    'error': 'Image descriptions disabled by user'
                }
            
            # Combine OCR text and image description
            combined_content = []
            
            # Add OCR text section
            combined_content.append("=== OCR EXTRACTED TEXT ===")
            combined_content.append(final_ocr_text)
            combined_content.append("")
            
            # Add image description section (if enabled)
            if enable_image_descriptions:
                combined_content.append("=== IMAGE DESCRIPTION (Gemma3) ===")
                if description_result['status'] == 'success':
                    combined_content.append(description_result['description'])
                    combined_content.append("")
                    combined_content.append(f"Model: {description_result.get('model_used', 'Unknown')}")
                    if 'response_metadata' in description_result:
                        metadata = description_result['response_metadata']
                        combined_content.append(f"Processing time: {metadata.get('total_duration', 0)/1000000000:.2f}s")
                else:
                    combined_content.append(f"Error getting image description: {description_result.get('error', 'Unknown error')}")
            else:
                combined_content.append("=== IMAGE DESCRIPTION ===")
                combined_content.append("Image descriptions disabled by user")
            
            return {
                'ocr_text': final_ocr_text,
                'image_description': description_result.get('description', ''),
                'description_status': description_result.get('status', 'error'),
                'description_error': description_result.get('error', ''),
                'combined_content': '\n'.join(combined_content),
                'status': 'success'
            }
            
        except Exception as e:
            logger.error(f"Error processing image file: {str(e)}")
            return {
                'ocr_text': f"Error reading image file: {str(e)}",
                'image_description': '',
                'description_status': 'error',
                'description_error': str(e),
                'combined_content': f"Error reading image file: {str(e)}",
                'status': 'error',
                'error': str(e)
            }
    
    def _process_powerpoint(self, file_path: str) -> str:
        """Process PowerPoint files"""
        try:
            prs = Presentation(file_path)
            content_parts = []
            
            for slide_num, slide in enumerate(prs.slides):
                content_parts.append(f"Slide {slide_num + 1}:")
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        content_parts.append(shape.text)
                    
                    # Handle tables specifically
                    elif hasattr(shape, "table"):
                        table = shape.table
                        table_content = []
                        
                        for row_idx, row in enumerate(table.rows):
                            row_cells = []
                            for cell in row.cells:
                                cell_text = cell.text.strip()
                                row_cells.append(cell_text if cell_text else "")
                            
                            # Only add non-empty rows
                            if any(cell for cell in row_cells):
                                table_content.append(" | ".join(row_cells))
                        
                        if table_content:
                            content_parts.append("Table:")
                            content_parts.extend(table_content)
                
                content_parts.append("\n" + "-"*30 + "\n")
            
            return "\n".join(content_parts)
            
        except Exception as e:
            logger.error(f"Error processing PowerPoint file: {str(e)}")
            return f"Error reading PowerPoint file: {str(e)}"
    
    def _process_word(self, file_path: str) -> str:
        """Process Word documents"""
        try:
            doc = Document(file_path)
            content_parts = []
            
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    content_parts.append(paragraph.text)
            
            return "\n".join(content_parts)
            
        except Exception as e:
            logger.error(f"Error processing Word file: {str(e)}")
            return f"Error reading Word file: {str(e)}"
    
    def batch_process(self, file_paths: List[str], enable_image_descriptions: bool = True) -> List[Dict[str, Any]]:
        """
        Process multiple files in batch
        
        Args:
            file_paths (List[str]): List of file paths to process
            enable_image_descriptions (bool): Whether to enable AI image descriptions for image files
            
        Returns:
            List of processing results
        """
        results = []
        for file_path in file_paths:
            result = self.process_file(file_path, enable_image_descriptions=enable_image_descriptions)
            results.append(result)
        
        return results
    
    def _ocr_with_preprocessing(self, image: np.ndarray, method: str) -> str:
        """Apply different preprocessing techniques for OCR"""
        try:
            gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
            
            if method == 'standard':
                # Standard preprocessing
                denoised = cv2.fastNlMeansDenoising(gray)
                _, thresh = cv2.threshold(denoised, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)
                config = '--oem 3 --psm 6'
                
            elif method == 'handwritten':
                # Enhanced preprocessing for handwritten text
                # Apply bilateral filter to reduce noise while preserving edges
                filtered = cv2.bilateralFilter(gray, 9, 75, 75)
                
                # Apply adaptive thresholding for better handwritten text
                thresh = cv2.adaptiveThreshold(
                    filtered, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C, 
                    cv2.THRESH_BINARY, 11, 2
                )
                
                # Morphological operations to clean up
                kernel = np.ones((1,1), np.uint8)
                thresh = cv2.morphologyEx(thresh, cv2.MORPH_CLOSE, kernel)
                
                # Use PSM 8 for single word recognition (better for handwritten)
                config = '--oem 3 --psm 8'
                
            elif method == 'table':
                # Table-specific preprocessing
                # Enhance contrast
                clahe = cv2.createCLAHE(clipLimit=2.0, tileGridSize=(8,8))
                enhanced = clahe.apply(gray)
                
                # Apply Gaussian blur to smooth
                blurred = cv2.GaussianBlur(enhanced, (3, 3), 0)
                
                # Apply thresholding
                _, thresh = cv2.threshold(blurred, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)
                
                # Use PSM 6 for uniform block of text (good for tables)
                config = '--oem 3 --psm 6'
                
            else:
                # Default to standard
                denoised = cv2.fastNlMeansDenoising(gray)
                _, thresh = cv2.threshold(denoised, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)
                config = '--oem 3 --psm 6'
            
            # Perform OCR
            text = pytesseract.image_to_string(thresh, config=config)
            return self._clean_ocr_text(text)
            
        except Exception as e:
            logger.error(f"Error in {method} OCR preprocessing: {str(e)}")
            return ""
    
    def _ocr_raw_image(self, image: np.ndarray) -> str:
        """OCR on raw image with different PSM modes"""
        try:
            results = []
            
            # Try different PSM modes for raw image
            psm_modes = [3, 4, 6, 8, 11, 12, 13]
            
            for psm in psm_modes:
                try:
                    config = f'--oem 3 --psm {psm}'
                    text = pytesseract.image_to_string(image, config=config)
                    cleaned = self._clean_ocr_text(text)
                    if cleaned.strip():
                        results.append(cleaned)
                except Exception as e:
                    logger.warning(f"PSM {psm} failed: {str(e)}")
                    continue
            
            # Return the result with the most content
            if results:
                return max(results, key=len)
            return ""
            
        except Exception as e:
            logger.error(f"Error in raw image OCR: {str(e)}")
            return ""
    
    def _combine_ocr_results(self, results: List[str]) -> str:
        """Combine multiple OCR results intelligently"""
        try:
            # Filter out empty results
            valid_results = [r for r in results if r.strip()]
            
            if not valid_results:
                return ""
            
            # If only one result, return it
            if len(valid_results) == 1:
                return valid_results[0]
            
            # For simple, clean text, prefer shorter results (less likely to have noise)
            # Check if any result looks like clean, simple text
            clean_results = []
            for result in valid_results:
                # Clean text indicators: shorter length, no special characters, readable format
                if (len(result.strip()) < 100 and 
                    not any(char in result for char in ['', '~', '']) and
                    len(result.split('\n')) <= 5):
                    clean_results.append(result)
            
            if clean_results:
                # Return the shortest clean result
                return min(clean_results, key=lambda x: len(x.strip()))
            
            # Try to detect table structure and combine accordingly
            combined = self._detect_and_parse_table_structure(valid_results)
            
            if combined:
                return combined
            
            # Fallback: return the result with the most readable content
            # Prefer results without garbled characters
            readable_results = [r for r in valid_results if not any(char in r for char in ['', '~', ''])]
            if readable_results:
                return max(readable_results, key=len)
            
            # Final fallback: return the longest result
            return max(valid_results, key=len)
            
        except Exception as e:
            logger.error(f"Error combining OCR results: {str(e)}")
            return results[0] if results else ""
    
    def _detect_and_parse_table_structure(self, results: List[str]) -> str:
        """Detect and parse table structure from OCR results"""
        try:
            # Look for table indicators
            table_indicators = ['visitors', 'log', 'book', 'sign', 'date', 'name', 'time', 'reason']
            
            best_result = None
            max_table_score = 0
            
            for result in results:
                score = 0
                result_lower = result.lower()
                
                # Score based on table indicators
                for indicator in table_indicators:
                    if indicator in result_lower:
                        score += 1
                
                # Score based on structured data patterns
                if re.search(r'\d{1,2}[/-]\d{1,2}[/-]\d{2,4}', result):  # Date pattern
                    score += 2
                if re.search(r'\d{1,2}:\d{2}\s*(am|pm)', result, re.IGNORECASE):  # Time pattern
                    score += 2
                if re.search(r'[A-Z][a-z]+\s+[A-Z][a-z]+', result):  # Name pattern
                    score += 1
                
                if score > max_table_score:
                    max_table_score = score
                    best_result = result
            
            if max_table_score >= 3:  # Threshold for table detection
                return self._format_table_output(best_result)
            
            return ""
            
        except Exception as e:
            logger.error(f"Error detecting table structure: {str(e)}")
            return ""
    
    def _format_table_output(self, text: str) -> str:
        """Format OCR text as a structured table"""
        try:
            lines = text.split('\n')
            formatted_lines = []
            
            # Look for header row
            header_found = False
            for i, line in enumerate(lines):
                line_clean = line.strip()
                if not line_clean:
                    continue
                
                # Check if this looks like a header
                if any(word in line_clean.lower() for word in ['visitors', 'log', 'book', 'sign']):
                    formatted_lines.append("VISITORS LOG BOOK")
                    formatted_lines.append("=" * 50)
                    formatted_lines.append("Columns: # | Date | Name | Reason for Visit | Time In | Time Out | Sign/Initials")
                    formatted_lines.append("-" * 50)
                    header_found = True
                    continue
                
                # Check if this looks like table data
                if re.search(r'\d+', line_clean):  # Contains numbers (likely row data)
                    # Try to parse as table row
                    parsed_row = self._parse_table_row(line_clean)
                    if parsed_row:
                        formatted_lines.append(parsed_row)
                else:
                    # Check if this looks like a name or reason
                    if re.search(r'[A-Z][a-z]+\s+[A-Z][a-z]+', line_clean):  # Name pattern
                        formatted_lines.append(f"Name/Reason: {line_clean}")
                    else:
                        # Regular text
                        formatted_lines.append(line_clean)
            
            return '\n'.join(formatted_lines)
            
        except Exception as e:
            logger.error(f"Error formatting table output: {str(e)}")
            return text
    
    def _parse_table_row(self, line: str) -> str:
        """Parse a line as a table row"""
        try:
            # Common table row patterns for visitors log
            # Pattern: number, date, name, reason, time_in, time_out, signature
            
            # Extract date pattern
            date_match = re.search(r'(\d{1,2}[/-]\d{1,2}[/-]\d{2,4})', line)
            date = date_match.group(1) if date_match else "N/A"
            
            # Extract time patterns - more flexible patterns
            time_patterns = [
                r'(\d{1,2}:\d{2}\s*(?:am|pm)?)',  # Standard time format
                r'(\d{1,2}\s*(?:am|pm))',         # Time without minutes
                r'(\d{1,2}:\d{2})',               # Time without am/pm
            ]
            
            time_matches = []
            for pattern in time_patterns:
                matches = re.findall(pattern, line, re.IGNORECASE)
                time_matches.extend(matches)
            
            time_in = time_matches[0] if len(time_matches) > 0 else "N/A"
            time_out = time_matches[1] if len(time_matches) > 1 else "N/A"
            
            # Extract name pattern (capitalized words) - more flexible
            name_patterns = [
                r'([A-Z][a-z]+\s+[A-Z][a-z]+)',  # First Last
                r'([A-Z][a-z]+\s+[A-Z][a-z]+\s+[A-Z][a-z]+)',  # First Middle Last
                r'([A-Z][a-z]+\s+[A-Z]\.\s*[A-Z][a-z]+)',  # First M. Last
            ]
            
            name = "N/A"
            for pattern in name_patterns:
                name_match = re.search(pattern, line)
                if name_match:
                    name = name_match.group(1)
                    break
            
            # Extract row number
            num_match = re.search(r'^(\d+)', line)
            row_num = num_match.group(1) if num_match else "N/A"
            
            # Extract reason for visit (text after name)
            reason_patterns = [
                r'(?:Sales\s+Demo|Meeting\s+\w+\s+\w+|Demo|Meeting)',
                r'([A-Z][a-z]+(?:\s+[A-Z][a-z]+)*)',
            ]
            
            reason = "N/A"
            for pattern in reason_patterns:
                reason_match = re.search(pattern, line, re.IGNORECASE)
                if reason_match:
                    reason = reason_match.group(0)
                    break
            
            # Format as structured row
            return f"Row {row_num}: Date={date}, Name={name}, Reason={reason}, Time In={time_in}, Time Out={time_out}"
            
        except Exception as e:
            logger.error(f"Error parsing table row: {str(e)}")
            return line
